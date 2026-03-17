import re
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

MAX_BULLETS = 5
MAX_CODE_LINES = 10
MAX_CODE_LINE_LENGTH = 55
MAX_DIAGRAM_STEPS = 6


# COLORS
BG_COLOR = RGBColor(12, 12, 18)
TITLE_BAR = RGBColor(25, 140, 255)
BULLET_COLOR = RGBColor(230, 230, 230)

CODE_BG = RGBColor(15, 15, 15)
CODE_TEXT = RGBColor(0, 255, 120)

DIAGRAM_BOX = RGBColor(60, 160, 255)
DIAGRAM_TEXT = RGBColor(255, 255, 255)

# -------------------------
# LAYOUT CONSTANTS
# -------------------------
MARGIN_X = Inches(0.7)
CONTENT_TOP = Inches(1.3)
CONTENT_HEIGHT = Inches(3.6)

LEFT_WIDTH = Inches(5.5)
RIGHT_WIDTH = Inches(6)

BOTTOM_TOP = Inches(5.4)
BOTTOM_HEIGHT = Inches(1.2)

MAX_BULLETS_BASE = 5
MAX_BULLETS_WITH_CODE = 4
MAX_BULLETS_WITH_DIAGRAM = 3

MAX_BULLET_CHARS = 90


# -------------------------
# PARSE AI OUTPUT
# -------------------------
def parse_ai_slides(ai_text):

    slides = []

    slide_pattern = r"\[SLIDE:\s*(.*?)\](.*?)(?=\[SLIDE:|\Z)"
    matches = re.findall(slide_pattern, ai_text, re.S)

    for title, content in matches:

        bullets = []
        codes = []
        diagrams = []

        codes = re.findall(r"\[CODE\](.*?)\[/CODE\]", content, re.S)
        diagrams = re.findall(r"\[DIAGRAM\](.*?)\[/DIAGRAM\]", content, re.S)

        for line in content.splitlines():
            line = line.strip()
            if line.startswith("•"):
                bullets.append(line.replace("•", "").strip())

        slides.append({
            "title": title.strip(),
            "bullets": bullets[:MAX_BULLETS],
            "code": [c.strip() for c in codes],
            "diagram": [d.strip() for d in diagrams]
        })

    return slides


# -------------------------
# BACKGROUND
# -------------------------
def add_background(slide, prs):

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height
    )

    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()


# -------------------------
# TITLE
# -------------------------
def add_title(slide, title, prs):

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        Inches(1)
    )

    fill = bar.fill
    fill.solid()
    fill.fore_color.rgb = TITLE_BAR
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.2),
        prs.slide_width - Inches(1),
        Inches(0.7)
    )

    tf = title_box.text_frame
    tf.text = title

    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


# -------------------------
# BULLETS
# -------------------------
def trim_bullets(bullets, max_items):
    trimmed = []

    for b in bullets[:max_items]:
        if len(b) > MAX_BULLET_CHARS:
            b = b[:MAX_BULLET_CHARS] + "…"
        trimmed.append(b)

    return trimmed


def add_bullets(slide, bullets, prs, has_code=False, has_diagram=False):

    # Dynamic bullet limits
    if has_code and has_diagram:
        max_bullets = 3
    elif has_code:
        max_bullets = MAX_BULLETS_WITH_CODE
    elif has_diagram:
        max_bullets = MAX_BULLETS_WITH_DIAGRAM
    else:
        max_bullets = MAX_BULLETS_BASE

    bullets = trim_bullets(bullets, max_bullets)

    box = slide.shapes.add_textbox(
        MARGIN_X,
        CONTENT_TOP,
        prs.slide_width - (MARGIN_X * 2),
        Inches(0.6 + len(bullets) * 0.35)
    )

    tf = box.text_frame
    tf.word_wrap = True

    if not bullets:
        return

    # Clear BEFORE adding anything
    tf.clear()

    # Add bullets cleanly
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]

        p.text = bullet
        p.level = 0

        p.font.size = Pt(20)
        p.font.color.rgb = BULLET_COLOR
        p.font.name = "Calibri"
        p.space_after = Pt(10)

    return len(bullets)


# -------------------------
# CODE BLOCK
# -------------------------
def add_code_block(slide, code):

    lines = code.split("\n")[:MAX_CODE_LINES]

    max_line_length = max((len(line) for line in lines), default=0)

    # -------------------------
    # Dynamic font scaling
    # -------------------------
    if max_line_length > 70:
        font_size = 12
    elif max_line_length > 55:
        font_size = 13
    else:
        font_size = 14

    # -------------------------
    # Dynamic width scaling
    # -------------------------
    if max_line_length > 70:
        box_width = RIGHT_WIDTH - Inches(0.5)
    else:
        box_width = RIGHT_WIDTH

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_X + LEFT_WIDTH + Inches(0.3),
        CONTENT_TOP,
        box_width,
        CONTENT_HEIGHT
    )

    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = CODE_BG

    box.line.color.rgb = RGBColor(0, 255, 120)

    tf = box.text_frame
    tf.word_wrap = False

    # Padding
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)

    # Trim long lines visually
    processed_lines = [
        line[:MAX_CODE_LINE_LENGTH] + ("…" if len(line) > MAX_CODE_LINE_LENGTH else "")
        for line in lines
    ]

    tf.text = processed_lines[0] if processed_lines else ""

    for line in processed_lines[1:]:
        p = tf.add_paragraph()
        p.text = line

    for p in tf.paragraphs:
        p.font.name = "Courier New"
        p.font.size = Pt(font_size)
        p.font.color.rgb = CODE_TEXT


# -------------------------
# HORIZONTAL DIAGRAM
# -------------------------
def add_diagram(slide, diagram_text, prs, top):

    steps = [s.strip() for s in diagram_text.split("→")][:MAX_DIAGRAM_STEPS]

    if not steps:
        return

    total_steps = len(steps)

    # -------------------------
    # TRUE available width
    # -------------------------
    slide_width = prs.slide_width
    usable_width = slide_width - (MARGIN_X * 2)

    spacing = Inches(0.15)

    step_width = (usable_width / total_steps) - spacing

    # Prevent boxes from becoming too wide
    max_step_width = Inches(2.2)
    step_width = min(step_width, max_step_width)

    # Center diagram
    total_diagram_width = total_steps * (step_width + spacing)
    start_x = (slide_width - total_diagram_width) / 2


    # -------------------------
    # Font scaling
    # -------------------------
    if total_steps >= 6:
        font_size = 11
    elif total_steps >= 5:
        font_size = 12
    else:
        font_size = 13

    for i, step in enumerate(steps):

        left = start_x + i * (step_width + spacing)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left,
            top,
            step_width,
            Inches(0.7)
        )

        fill = box.fill
        fill.solid()
        fill.fore_color.rgb = DIAGRAM_BOX

        box.line.color.rgb = RGBColor(0, 255, 200)

        tf = box.text_frame
        tf.word_wrap = True
        tf.text = step

        p = tf.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = DIAGRAM_TEXT
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

def add_code_block_full_width(slide, code, prs, top):

    lines = code.split("\n")[:MAX_CODE_LINES]

    max_line_length = max((len(line) for line in lines), default=0)

    # Font scaling
    if max_line_length > 70:
        font_size = 12
    elif max_line_length > 55:
        font_size = 13
    else:
        font_size = 14

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN_X,
        top,
        prs.slide_width - (MARGIN_X * 2),
        Inches(1.6)
    )

    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = CODE_BG

    box.line.color.rgb = RGBColor(0, 255, 120)

    tf = box.text_frame
    tf.word_wrap = False

    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    lines = [
        line[:MAX_CODE_LINE_LENGTH] + ("…" if len(line) > MAX_CODE_LINE_LENGTH else "")
        for line in lines
    ]

    tf.text = lines[0] if lines else ""

    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line

    for p in tf.paragraphs:
        p.font.name = "Courier New"
        p.font.size = Pt(font_size)
        p.font.color.rgb = CODE_TEXT


# -------------------------
# GENERATE PPT
# -------------------------
def generate_ppt(slides):

    prs = Presentation()

    for slide_data in slides:

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        add_background(slide, prs)
        add_title(slide, slide_data["title"], prs)

        # -------------------------
        # Layout Switching Logic
        # -------------------------
        has_code = bool(slide_data["code"])
        has_diagram = bool(slide_data["diagram"])

        # -------------------------
        # Dynamic layout calculation
        # -------------------------
        actual_bullets = add_bullets(
            slide,
            slide_data["bullets"],
            prs,
            has_code=has_code,
            has_diagram=has_diagram
        )
        bullet_height = 0.6 + actual_bullets * 0.35
        code_top = CONTENT_TOP + Inches(bullet_height)

        # Code
        if has_code:
            add_code_block_full_width(
                slide,
                slide_data["code"][0],
                prs,
                code_top
            )

        # Diagram
        if has_diagram:
            if has_code:
                diagram_top = code_top + Inches(2.0)
            else:
                diagram_top = CONTENT_TOP + Inches(bullet_height + 0.5)

            add_diagram(
                slide,
                slide_data["diagram"][0],
                prs,
                diagram_top
            )

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)

    return ppt_stream


# -------------------------
# ENTRY FUNCTION
# -------------------------
def create_slides_from_ai(ai_text, project_name="project"):

    slides = parse_ai_slides(ai_text)

    if not slides:
        raise ValueError("No slides detected")

    return generate_ppt(slides)